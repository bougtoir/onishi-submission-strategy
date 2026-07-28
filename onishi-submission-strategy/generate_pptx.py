#!/usr/bin/env python3
"""
Generate PPTX files for ONISHI framework figures.
- Flow diagrams / conceptual diagrams → editable PowerPoint shapes
- Code-output charts (radar) → embedded as images
- One figure per slide
- English and Japanese versions
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import os

# =========================================================
# Colors
# =========================================================
C_LINKO = RGBColor(0x21, 0x96, 0xF3)   # Blue
C_IONE  = RGBColor(0xFF, 0x98, 0x00)   # Orange
C_KOTHA = RGBColor(0x4C, 0xAF, 0x50)   # Green
C_ALL3  = RGBColor(0x9C, 0x27, 0xB0)   # Purple
C_DARK  = RGBColor(0x37, 0x47, 0x4F)   # Dark gray
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x33, 0x33, 0x33)
C_GRAY  = RGBColor(0x78, 0x90, 0x9C)
C_LGRAY = RGBColor(0xE0, 0xE0, 0xE0)
C_ARROW = RGBColor(0x61, 0x61, 0x61)

# Light backgrounds for step boxes
C_LIGHT_BLUE   = RGBColor(0xE3, 0xF2, 0xFD)
C_LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
C_LIGHT_GREEN  = RGBColor(0xE8, 0xF5, 0xE9)
C_LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)

RADAR_IMG = "/home/ubuntu/report_figures/fig5_radar.png"

# Slide dimensions: widescreen 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rounded_box(slide, left, top, width, height, fill_color, text="",
                    font_size=11, font_color=C_WHITE, bold=True, border_color=None,
                    border_width=Pt(0), alignment=PP_ALIGN.CENTER):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    # Adjust corner rounding
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = font_color
    tf.paragraphs[0].font.bold = bold

    if text:
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
                tf.paragraphs[0].font.size = Pt(font_size)
                tf.paragraphs[0].font.color.rgb = font_color
                tf.paragraphs[0].font.bold = bold
                tf.paragraphs[0].alignment = alignment
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(font_size)
                p.font.color.rgb = font_color
                p.font.bold = bold
                p.alignment = alignment

    return shape


def add_rich_box(slide, left, top, width, height, fill_color, paragraphs_spec,
                 border_color=None, border_width=Pt(0), v_anchor=MSO_ANCHOR.TOP):
    """Add a rounded rectangle with multiple styled paragraphs.
    paragraphs_spec: list of dicts with keys:
        text, font_size, font_color, bold, italic, alignment, space_before, space_after
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = v_anchor

    for i, spec in enumerate(paragraphs_spec):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = spec.get("text", "")
        p.font.size = Pt(spec.get("font_size", 11))
        p.font.color.rgb = spec.get("font_color", C_WHITE)
        p.font.bold = spec.get("bold", False)
        p.font.italic = spec.get("italic", False)
        p.alignment = spec.get("alignment", PP_ALIGN.CENTER)
        if "space_before" in spec:
            p.space_before = spec["space_before"]
        if "space_after" in spec:
            p.space_after = spec["space_after"]

    return shape


def add_textbox(slide, left, top, width, height, text, font_size=11,
                font_color=C_BLACK, bold=False, alignment=PP_ALIGN.CENTER,
                italic=False):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_arrow_down(slide, x, y_start, y_end, color=C_ARROW):
    """Add a downward arrow using a thin triangle shape."""
    x_int = int(x)
    y_start_int = int(y_start)
    y_end_int = int(y_end)
    length = y_end_int - y_start_int
    # Use a line shape with arrow
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, x_int - Inches(0.1), y_start_int, Inches(0.2), length
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_right(slide, x_start, y, x_end, color=C_ARROW):
    """Add a rightward arrow shape."""
    x_start_int = int(x_start)
    y_int = int(y)
    x_end_int = int(x_end)
    length = x_end_int - x_start_int
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x_start_int, y_int - Inches(0.1), length, Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# =========================================================
# Text content for English and Japanese
# =========================================================
def get_texts(lang="en"):
    if lang == "en":
        return {
            # Slide 1: Overview
            "s1_title": "Three Analytical Frameworks: LINKO, IONE, KOTHA",
            "s1_subtitle": "Components of the ONISHI Integrated Framework\n(Optimal Normalization, Incoherence Stratification, and Harmonized Integration)",
            "linko_name": "LINKO",
            "linko_full": "Latent Information Normalization\nfor Key Outcomes",
            "linko_items": ["ICR (Information Contribution Ratio)",
                            "Meta-analysis validity diagnostic",
                            "Prism Forest Plot visualization",
                            "ICR-weighted pooling"],
            "linko_target": "Target: Meta-analysis level",
            "ione_name": "IONE",
            "ione_full": "Incoherence-Oriented Neutralization\nand Extraction",
            "ione_items": ["Coherent subgroup detection",
                           "C1 coherence index",
                           "Simpson's paradox resolution",
                           "Hidden structure identification"],
            "ione_target": "Target: Individual study level",
            "kotha_name": "KOTHA",
            "kotha_full": "Knowledge-driven Observational-Trial\nHarmonization Approach",
            "kotha_items": ["Module K: Counterfactual power sim.",
                            "Module T: Bayesian integration",
                            "Module H: Guideline interpreter",
                            "RCT-observational harmonization"],
            "kotha_target": "Target: Evidence synthesis level",
            "common_data": "Common Data Sources",
            "common_data_detail": "RCT published data (Table 1)  |  Individual Patient Data (IPD)  |  Observational cohort data",

            # Slide 2: LINKO + IONE
            "s2_title": "Combination 1: LINKO + IONE",
            "s2_subtitle": "ICR-guided Subpopulation Analysis",
            "s2_steps": [
                ("Step 1", "LINKO identifies anomalous ICR across studies in meta-analysis"),
                ("Step 2", "IONE decomposes each study into coherent subgroups"),
                ("Step 3", "ICR recalculated per subgroup reveals hidden heterogeneity sources"),
                ("Step 4", "Subgroup-specific ICR-weighted meta-analysis"),
            ],
            "s2_gain": "Gain: Root cause identification of heterogeneity",

            # Slide 3: LINKO + KOTHA
            "s3_title": "Combination 2: LINKO + KOTHA",
            "s3_subtitle": "Information-Aware Evidence Synthesis",
            "s3_steps": [
                ("Step 1", "LINKO quantifies information contribution of endpoints"),
                ("Step 2", "ICR feeds into KOTHA Module K for refined power simulation"),
                ("Step 3", "Module T uses ICR as Bayesian prior weight for integration"),
                ("Step 4", "Module H reports ICR alongside OIS and TSA assessments"),
            ],
            "s3_gain": "Gain: Quantitative precision in evidence synthesis",

            # Slide 4: IONE + KOTHA
            "s4_title": "Combination 3: IONE + KOTHA",
            "s4_subtitle": "Population-Harmonized Evidence Integration",
            "s4_steps": [
                ("Step 1", "IONE detects incoherent subgroups in observational data"),
                ("Step 2", "Subgroup risk profiles inform KOTHA Module K simulation"),
                ("Step 3", "Module T integrates subgroup-specific effects with RCT data"),
                ("Step 4", "Module H assesses representativeness using IONE's C1 index"),
            ],
            "s4_gain": "Gain: External validity through population harmonization",

            # Slide 5: LINKO + IONE + KOTHA
            "s5_title": "Combination 4: LINKO + IONE + KOTHA  =  ONISHI",
            "s5_subtitle": "Comprehensive Evidence Assessment Pipeline\n(Optimal Normalization, Incoherence Stratification, and Harmonized Integration)",
            "s5_steps": [
                ("Phase 1\nIONE", "Population Decomposition\nDetect incoherent populations → Extract coherent subgroups → Compute C1 index"),
                ("Phase 2\nLINKO", "Information Quantification\nCompute ICR per subgroup → Assess ICR discrepancy → Prism Forest Plot"),
                ("Phase 3\nKOTHA", "Evidence Harmonization\nModule K: Power simulation → Module T: Bayesian synthesis → Module H: Guidelines"),
                ("Output", "Integrated Decision Support\nSubgroup-specific recommendations with full uncertainty quantification"),
            ],
            "s5_gain": "Gain: Comprehensive multi-level evidence assessment",

            # Slide 6: Synergy Matrix
            "s6_title": "Synergy Matrix: What Each Combination Enables",
            "s6_headers": ["Combination", "New Capability", "Primary Gain"],
            "s6_rows": [
                ("LINKO + IONE", "Subgroup-level ICR analysis reveals\nwhy meta-analysis heterogeneity arises", "Root cause\nidentification"),
                ("LINKO + KOTHA", "ICR-informed power simulation and\nBayesian weighting improve precision", "Quantitative\nprecision"),
                ("IONE + KOTHA", "Population decomposition feeds\ncounterfactual simulation & representativeness", "External\nvalidity"),
                ("ONISHI\n(All three)", "Full pipeline: detect heterogeneity →\nquantify information → harmonize evidence", "Comprehensive\nassessment"),
            ],

            # Slide 7: Radar chart title
            "s7_title": "Capability Profile by Combination",
            "s7_subtitle": "Comparative radar chart across 6 dimensions",
        }
    else:  # Japanese
        return {
            # Slide 1: Overview
            "s1_title": "3つの分析フレームワーク: LINKO, IONE, KOTHA",
            "s1_subtitle": "ONISHI 統合フレームワークの構成要素\n(Optimal Normalization, Incoherence Stratification, and Harmonized Integration)",
            "linko_name": "LINKO",
            "linko_full": "Latent Information Normalization\nfor Key Outcomes",
            "linko_items": ["ICR（情報寄与率）",
                            "メタ解析の妥当性診断",
                            "Prism Forest Plot 可視化",
                            "ICR加重プーリング"],
            "linko_target": "対象: メタ解析レベル",
            "ione_name": "IONE",
            "ione_full": "Incoherence-Oriented Neutralization\nand Extraction",
            "ione_items": ["コヒーレント部分集団の検出",
                           "C1 コヒーレンス指標",
                           "シンプソンのパラドックス解消",
                           "隠れた構造の同定"],
            "ione_target": "対象: 個別研究レベル",
            "kotha_name": "KOTHA",
            "kotha_full": "Knowledge-driven Observational-Trial\nHarmonization Approach",
            "kotha_items": ["Module K: 反事実的検出力シミュレーション",
                            "Module T: ベイズ統合",
                            "Module H: ガイドライン解釈器",
                            "RCT-観察研究の調和"],
            "kotha_target": "対象: エビデンス統合レベル",
            "common_data": "共通データソース",
            "common_data_detail": "RCT公表データ（Table 1）｜ 個人患者データ（IPD）｜ 観察コホートデータ",

            # Slide 2: LINKO + IONE
            "s2_title": "組み合わせ 1: LINKO + IONE",
            "s2_subtitle": "ICR指導型 部分集団解析",
            "s2_steps": [
                ("Step 1", "LINKOがメタ解析中の異常なICRを同定"),
                ("Step 2", "IONEが各研究をコヒーレントな部分集団に分解"),
                ("Step 3", "部分集団ごとにICRを再計算し、異質性の隠れた原因を特定"),
                ("Step 4", "部分集団別ICR加重メタ解析を実施"),
            ],
            "s2_gain": "獲得能力: 異質性の根本原因の特定",

            # Slide 3: LINKO + KOTHA
            "s3_title": "組み合わせ 2: LINKO + KOTHA",
            "s3_subtitle": "情報構造を反映したエビデンス統合",
            "s3_steps": [
                ("Step 1", "LINKOがエンドポイントの情報寄与を定量化"),
                ("Step 2", "ICRがKOTHA Module Kの精密検出力シミュレーションに入力"),
                ("Step 3", "Module TがICRをベイズ事前重みとして統合に使用"),
                ("Step 4", "Module HがOIS・TSA評価とともにICRを報告"),
            ],
            "s3_gain": "獲得能力: エビデンス統合の定量的精度向上",

            # Slide 4: IONE + KOTHA
            "s4_title": "組み合わせ 3: IONE + KOTHA",
            "s4_subtitle": "集団構造を反映したエビデンス統合",
            "s4_steps": [
                ("Step 1", "IONEが観察データ中の非一貫的部分集団を検出"),
                ("Step 2", "部分集団のリスクプロファイルがKOTHA Module Kに入力"),
                ("Step 3", "Module Tが部分集団特異的効果とRCTデータを統合"),
                ("Step 4", "Module HがIONEのC1指標で代表性を評価"),
            ],
            "s4_gain": "獲得能力: 集団調和による外的妥当性の向上",

            # Slide 5: LINKO + IONE + KOTHA
            "s5_title": "組み合わせ 4: LINKO + IONE + KOTHA = ONISHI",
            "s5_subtitle": "包括的エビデンス評価パイプライン\n(Optimal Normalization, Incoherence Stratification, and Harmonized Integration)",
            "s5_steps": [
                ("Phase 1\nIONE", "集団分解\n非一貫的集団を検出 → コヒーレント部分集団を抽出 → C1指標を算出"),
                ("Phase 2\nLINKO", "情報定量化\n部分集団ごとにICRを算出 → ICR乖離を評価 → Prism Forest Plot"),
                ("Phase 3\nKOTHA", "エビデンス調和\nModule K: 検出力シミュレーション → Module T: ベイズ統合 → Module H: ガイドライン"),
                ("出力", "統合的意思決定支援\n部分集団別推奨事項（完全な不確実性定量化付き）"),
            ],
            "s5_gain": "獲得能力: 多層的な包括的エビデンス評価",

            # Slide 6: Synergy Matrix
            "s6_title": "シナジーマトリックス: 各組み合わせで可能になること",
            "s6_headers": ["組み合わせ", "新たに可能になる能力", "主な獲得"],
            "s6_rows": [
                ("LINKO + IONE", "部分集団レベルのICR解析により\nメタ解析の異質性の原因を解明", "根本原因\nの特定"),
                ("LINKO + KOTHA", "ICR情報に基づく検出力シミュレーション\nとベイズ加重で精度向上", "定量的\n精度"),
                ("IONE + KOTHA", "集団分解が反事実シミュレーション\nと代表性評価に貢献", "外的\n妥当性"),
                ("ONISHI\n（3手法統合）", "完全パイプライン: 異質性検出 →\n情報定量化 → エビデンス調和", "包括的\n評価"),
            ],

            # Slide 7: Radar chart title
            "s7_title": "組み合わせ別 能力プロファイル",
            "s7_subtitle": "6次元のレーダーチャートによる比較",
        }


# =========================================================
# Slide builders
# =========================================================

def build_slide1_overview(prs, t):
    """Three Frameworks Overview - editable shapes, all text inside shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                t["s1_title"], font_size=24, bold=True, font_color=C_BLACK)
    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
                t["s1_subtitle"], font_size=12, font_color=C_GRAY, italic=True)

    # Three framework boxes
    box_w = Inches(3.6)
    box_h = Inches(3.2)
    box_top = Inches(1.6)
    gap = Inches(0.4)
    start_x = Inches(0.8)

    frameworks = [
        (t["linko_name"], t["linko_full"], t["linko_items"], t["linko_target"], C_LINKO),
        (t["ione_name"], t["ione_full"], t["ione_items"], t["ione_target"], C_IONE),
        (t["kotha_name"], t["kotha_full"], t["kotha_items"], t["kotha_target"], C_KOTHA),
    ]

    for i, (name, full_name, items, target, color) in enumerate(frameworks):
        x = start_x + i * (box_w + gap)

        # Build all text as paragraphs inside one shape
        paras = [
            {"text": name, "font_size": 18, "font_color": C_WHITE, "bold": True,
             "space_after": Pt(2)},
            {"text": full_name.replace('\n', ' '), "font_size": 9, "font_color": RGBColor(0xE0, 0xE0, 0xE0),
             "italic": True, "space_after": Pt(6)},
        ]
        for item in items:
            paras.append({"text": f"  {item}", "font_size": 10, "font_color": C_WHITE,
                          "alignment": PP_ALIGN.LEFT, "space_before": Pt(2)})

        add_rich_box(slide, x, box_top, box_w, box_h, color, paras,
                     v_anchor=MSO_ANCHOR.TOP)

        # Target level label (below the box, separate textbox is fine - no overlap)
        add_textbox(slide, x, box_top + box_h + Inches(0.05), box_w, Inches(0.3),
                    target, font_size=9, font_color=color, italic=True)

    # Common data layer - all text inside one shape
    data_top = Inches(5.4)
    add_rich_box(slide, Inches(1.5), data_top, Inches(10), Inches(1.2), C_GRAY, [
        {"text": t["common_data"], "font_size": 14, "font_color": C_WHITE, "bold": True,
         "space_after": Pt(4)},
        {"text": t["common_data_detail"], "font_size": 10, "font_color": C_WHITE},
    ], v_anchor=MSO_ANCHOR.MIDDLE)


def build_slide_combination(prs, t, prefix, colors_pair, step_colors):
    """Build a combination slide (slides 2-4) with editable flow.
    All text is embedded directly inside shapes - no overlapping text boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
                t[f"{prefix}_title"], font_size=22, bold=True, font_color=C_BLACK)
    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(0.65), Inches(9), Inches(0.4),
                t[f"{prefix}_subtitle"], font_size=14, font_color=C_GRAY, italic=True)

    # Color indicator bars (top right, no text overlap with title)
    bar_y = Inches(0.15)
    if len(colors_pair) == 2:
        add_rounded_box(slide, Inches(10), bar_y, Inches(1.5), Inches(0.4), colors_pair[0])
        add_rounded_box(slide, Inches(11.5), bar_y, Inches(1.5), Inches(0.4), colors_pair[1])
    else:
        add_rounded_box(slide, Inches(9.5), bar_y, Inches(1.1), Inches(0.4), colors_pair[0])
        add_rounded_box(slide, Inches(10.6), bar_y, Inches(1.1), Inches(0.4), colors_pair[1])
        add_rounded_box(slide, Inches(11.7), bar_y, Inches(1.1), Inches(0.4), colors_pair[2])

    # Steps as flow boxes - text embedded directly in each box
    steps = t[f"{prefix}_steps"]
    box_w = Inches(10)
    box_h = Inches(0.95)
    start_x = Inches(1.5)
    start_y = Inches(1.3)
    gap_y = Inches(0.25)

    border_colors_list = [C_LINKO, C_IONE, C_KOTHA, C_ALL3]

    for i, (step_label, step_desc) in enumerate(steps):
        y = start_y + i * (box_h + gap_y)
        bc = border_colors_list[i] if i < len(border_colors_list) else C_ALL3

        # Step box with label + description inside
        add_rich_box(slide, start_x, y, box_w, box_h,
                     step_colors[i], [
                         {"text": f"{step_label}:  {step_desc}",
                          "font_size": 12, "font_color": C_BLACK,
                          "alignment": PP_ALIGN.LEFT},
                     ],
                     border_color=bc, border_width=Pt(2),
                     v_anchor=MSO_ANCHOR.MIDDLE)

        # Arrow between steps
        if i < len(steps) - 1:
            arrow_x = start_x + box_w / 2
            arrow_y_start = y + box_h
            arrow_y_end = y + box_h + gap_y
            add_arrow_down(slide, arrow_x, arrow_y_start, arrow_y_end, C_ARROW)

    # Gain box at bottom - text inside
    gain_y = start_y + len(steps) * (box_h + gap_y) + Inches(0.1)
    gain_color = colors_pair[0] if len(colors_pair) == 2 else C_ALL3
    add_rounded_box(slide, Inches(3), gain_y, Inches(7), Inches(0.6),
                    gain_color, t[f"{prefix}_gain"], font_size=13, font_color=C_WHITE)


def build_slide5_pipeline(prs, t):
    """Full ONISHI pipeline - all text inside shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_textbox(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                t["s5_title"], font_size=22, bold=True, font_color=C_BLACK)
    # Subtitle
    add_textbox(slide, Inches(0.3), Inches(0.6), Inches(12.5), Inches(0.5),
                t["s5_subtitle"], font_size=11, font_color=C_GRAY, italic=True)

    # Pipeline phases as horizontal flow
    phase_colors = [C_IONE, C_LINKO, C_KOTHA, C_ALL3]
    phase_w = Inches(2.8)
    phase_h = Inches(3.8)
    start_x = Inches(0.5)
    phase_y = Inches(1.5)
    gap_x = Inches(0.35)

    steps = t["s5_steps"]

    for i, (phase_label, phase_desc) in enumerate(steps):
        x = start_x + i * (phase_w + gap_x)

        # Phase box with label + description inside as paragraphs
        label_lines = phase_label.split("\n")
        desc_lines = phase_desc.split("\n")
        paras = []
        for ll in label_lines:
            paras.append({"text": ll, "font_size": 14, "font_color": C_WHITE,
                          "bold": True, "space_after": Pt(2)})
        paras.append({"text": "", "font_size": 6, "font_color": C_WHITE})  # spacer
        for dl in desc_lines:
            paras.append({"text": dl, "font_size": 11, "font_color": C_WHITE,
                          "alignment": PP_ALIGN.LEFT, "space_before": Pt(2)})

        add_rich_box(slide, x, phase_y, phase_w, phase_h,
                     phase_colors[i], paras, v_anchor=MSO_ANCHOR.TOP)

        # Arrow between phases
        if i < len(steps) - 1:
            arrow_x_start = x + phase_w
            arrow_x_end = x + phase_w + gap_x
            arrow_y = phase_y + phase_h / 2
            add_arrow_right(slide, arrow_x_start, arrow_y, arrow_x_end, C_ARROW)

    # Gain box at bottom - text inside
    gain_y = phase_y + phase_h + Inches(0.3)
    add_rounded_box(slide, Inches(2.5), gain_y, Inches(8), Inches(0.6),
                    C_ALL3, t["s5_gain"], font_size=14, font_color=C_WHITE)


def build_slide6_synergy(prs, t):
    """Synergy Matrix - editable table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                t["s6_title"], font_size=22, bold=True, font_color=C_BLACK)

    # Build as a table
    rows_data = t["s6_rows"]
    n_rows = len(rows_data) + 1  # +1 for header
    n_cols = 3
    table_left = Inches(0.8)
    table_top = Inches(1.2)
    table_w = Inches(11.5)
    table_h = Inches(5.5)

    table = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(6.0)
    table.columns[2].width = Inches(3.0)

    # Header row
    headers = t["s6_headers"]
    combo_colors = [C_LINKO, C_LINKO, C_IONE, C_ALL3]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = C_WHITE
        p.font.size = Pt(13)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, (combo, capability, gain) in enumerate(rows_data):
        row_idx = i + 1
        row_data = [combo, capability, gain]
        for j, val in enumerate(row_data):
            cell = table.cell(row_idx, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = combo_colors[i]
            elif j == 2:
                p.font.bold = True
                p.font.color.rgb = combo_colors[i]
            else:
                p.font.color.rgb = C_BLACK
            p.alignment = PP_ALIGN.CENTER
            # Alternate row background
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE


def build_slide7_radar(prs, t):
    """Radar chart - embedded as image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                t["s7_title"], font_size=22, bold=True, font_color=C_BLACK)
    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                t["s7_subtitle"], font_size=13, font_color=C_GRAY, italic=True)

    # Embed radar chart image
    if os.path.exists(RADAR_IMG):
        img_left = Inches(3.0)
        img_top = Inches(1.3)
        img_w = Inches(7.0)
        slide.shapes.add_picture(RADAR_IMG, img_left, img_top, width=img_w)


# =========================================================
# Main generator
# =========================================================
def generate_pptx(lang="en"):
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    t = get_texts(lang)

    # Slide 1: Overview
    build_slide1_overview(prs, t)

    # Slide 2: LINKO + IONE
    build_slide_combination(prs, t, "s2",
                            [C_LINKO, C_IONE],
                            [C_LIGHT_BLUE, C_LIGHT_ORANGE, C_LIGHT_GREEN, C_LIGHT_PURPLE])

    # Slide 3: LINKO + KOTHA
    build_slide_combination(prs, t, "s3",
                            [C_LINKO, C_KOTHA],
                            [C_LIGHT_BLUE, C_LIGHT_GREEN, C_LIGHT_BLUE, C_LIGHT_GREEN])

    # Slide 4: IONE + KOTHA
    build_slide_combination(prs, t, "s4",
                            [C_IONE, C_KOTHA],
                            [C_LIGHT_ORANGE, C_LIGHT_GREEN, C_LIGHT_ORANGE, C_LIGHT_GREEN])

    # Slide 5: Full ONISHI pipeline
    build_slide5_pipeline(prs, t)

    # Slide 6: Synergy matrix
    build_slide6_synergy(prs, t)

    # Slide 7: Radar chart (image)
    build_slide7_radar(prs, t)

    suffix = "EN" if lang == "en" else "JP"
    out_path = f"/home/ubuntu/ONISHI_figures_{suffix}.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    en_path = generate_pptx("en")
    jp_path = generate_pptx("ja")
    print(f"\nDone!\n  English: {en_path}\n  Japanese: {jp_path}")
