# -*- coding: utf-8 -*-
"""Build the complete AJE submission package for the ONISHI integration paper.

Outputs (into submission/build/):
  - ONISHI_manuscript.docx         main manuscript (submission format:
                                   figures NOT embedded; tables + legends at end)
  - ONISHI_manuscript_inline.docx  review copy with figures/tables placed inline
  - ONISHI_tables.docx             editable tables only
  - ONISHI_figures.pptx            editable figures, one per slide
  - ONISHI_cover_letter.docx       cover letter (discloses AI assistance)
  - ONISHI_highlights_keywords.docx
  - ONISHI_reviewer_suggestions.docx
  - ONISHI_STROBE_checklist.docx
  - figures/Figure_1..6.(png|tif)
  - ONISHI_submission_package.zip  everything above
"""

import os
import re
import json
import math
import shutil
import zipfile
import datetime

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.text import PP_ALIGN

import content as C

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
RESULTS = os.path.join(ROOT, "integration_analysis", "results.json")
INTEG_FIGDIR = os.path.join(ROOT, "integration_analysis", "figures")
BUILD = os.path.join(HERE, "build")
FIGOUT = os.path.join(BUILD, "figures")

FONT = "Times New Roman"

# ---------------------------------------------------------------------------
# Numbers
# ---------------------------------------------------------------------------
with open(RESULTS) as f:
    R = json.load(f)


def ci_from_log(logor, se):
    lo = math.exp(logor - 1.96 * se)
    hi = math.exp(logor + 1.96 * se)
    return lo, hi


def pct(x, d=1):
    return f"{100*x:.{d}f}%"


def sci(x):
    """Scientific notation like 2.3x10^-5 using unicode superscript exponent."""
    exp = int(math.floor(math.log10(abs(x))))
    mant = x / (10 ** exp)
    sup = str(exp).translate(str.maketrans("-0123456789",
                                            "\u207b\u2070\u00b9\u00b2\u00b3"
                                            "\u2074\u2075\u2076\u2077\u2078\u2079"))
    return f"{mant:.1f}\u00d710{sup}"


ov = R["overall_effect"]
ov_lo, ov_hi = ci_from_log(ov["logOR"], ov["se"])
strata = R["ione_stratum_effects"]
p1 = R["pattern1_linko_ione"]["table"]
p2 = R["pattern2_linko_kotha"]
p3 = R["pattern3_ione_kotha"]
p4 = R["pattern4_onishi_full"]

icr_vals = [row["icr_pca_reg"] for row in p1]
icr_med = sorted(icr_vals)[1:3]
icr_med = sum(icr_med) / 2

# Power-prior harmonized OR range: the naive and ICR-guided pooled odds ratios
# bound the harmonization; both derive from results.json (no hardcoding).
pp_ors = sorted([p4["naive_pooled_OR"], p4["icr_pooled_OR"]])

# Between-country ICR weights (near-uniform); range derived from results.json.
_bs_icr = [row["icr_pca_reg"] for row in p2["table"]]

VALUES = {
    "N": f"{R['dataset']['N']:,}",
    "N_RAND": f"{R['dataset']['n_randomized']:,}",
    "ER": pct(R["dataset"]["event_rate"]),
    "OR_overall": f"{ov['OR']:.2f}",
    "CI_overall": f"{ov_lo:.2f}\u2013{ov_hi:.2f}",
    "ER_s0": pct(strata[0]["event_rate"]),
    "ER_s3": pct(strata[3]["event_rate"]),
    "C1": f"{R['ione']['c1_effect']:.2f}",
    "W": f"{R['ione']['W']:.2f}",
    "ICR_s0": sci(p1[0]["icr_pca_reg"]),
    "ICR_s1": sci(p1[1]["icr_pca_reg"]),
    "ICR_s2": sci(p1[2]["icr_pca_reg"]),
    "ICR_s3": sci(p1[3]["icr_pca_reg"]),
    "ICR_med": sci(icr_med),
    "OR_iv": f"{p2['pooled_iv_OR']:.2f}",
    "SE_iv": f"{p2['pooled_iv_se']:.3f}",
    "I2_p2": f"{p2['I2']:.0f}%",
    "OR_icr": f"{p2['pooled_icr_OR']:.2f}",
    "SE_icr": f"{p2['pooled_icr_se']:.3f}",
    "SE_unit": f"{p2['pooled_unit_se']:.3f}",
    "pow_s0": pct(p3["table"][0]["power_at_pooled_OR"], 0),
    "pow_s3": pct(p3["table"][3]["power_at_pooled_OR"], 0),
    "OIS": f"{p3['ois_events']:,.0f}",
    "OBS": f"{p3['total_events']:,}",
    "INFOFRAC": pct(p3["info_fraction"], 0),
    "TSAZ": f"\u2212{abs(p4['tsa_final_z']):.2f}",
    "OR_pp_lo": f"{pp_ors[0]:.2f}",
    "OR_pp_hi": f"{pp_ors[1]:.2f}",
    "PB_naive": f"{p4['naive_p_benefit']:.2f}",
    "PB_icr": f"{p4['icr_p_benefit']:.2f}",
    "ADD_EV": f"{p3['additional_events_needed']:,.0f}",
    # rounded to the nearest 100 additional patients for the prose; derived from
    # results.
    "ADD_N": f"{round(p3['additional_n_overall'], -2):,.0f}",
    # per-stratum additional patients to reach the OIS (nearest 100); enrichment
    # is far more efficient in the highest-risk stratum than the lowest.
    "ADD_N_S3": f"{round(p3['additional_n_by_stratum'][3]['additional_n'], -2):,.0f}",
    "ADD_N_S0": f"{round(p3['additional_n_by_stratum'][0]['additional_n'], -2):,.0f}",
    "ICR_BS_LO": f"{min(_bs_icr):.3f}",
    "ICR_BS_HI": f"{max(_bs_icr):.3f}",
}

VAL_RE = re.compile(r"\{([A-Za-z_][A-Za-z0-9_]*)\}")


def fill(text):
    return VAL_RE.sub(lambda m: VALUES.get(m.group(1), m.group(0)), text)


CIT_RE = re.compile(r"(\{[^}]+\})")


def add_runs(paragraph, text, size=11, bold=False, italic=False):
    """Add text to paragraph, rendering {..} tokens as superscript citations."""
    for part in CIT_RE.split(text):
        if not part:
            continue
        if part.startswith("{") and part.endswith("}"):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
        else:
            run = paragraph.add_run(part)
        run.font.name = FONT
        run.font.size = Pt(size)
        run.bold = bold
        run.italic = italic


def set_base_style(doc):
    st = doc.styles["Normal"]
    st.font.name = FONT
    st.font.size = Pt(11)


def add_line_numbers(section):
    sectPr = section._sectPr
    ln = OxmlElement("w:lnNumType")
    ln.set(qn("w:countBy"), "1")
    ln.set(qn("w:restart"), "continuous")
    ln.set(qn("w:distance"), "360")
    sectPr.append(ln)


def add_page_numbers(section):
    footer = section.footer
    p = footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    fldb = OxmlElement("w:fldChar"); fldb.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText"); instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    flde = OxmlElement("w:fldChar"); flde.set(qn("w:fldCharType"), "end")
    run._r.append(fldb); run._r.append(instr); run._r.append(flde)


def heading(doc, text, size=12, space_before=14, space_after=6):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    r = p.add_run(text)
    r.bold = True
    r.font.name = FONT
    r.font.size = Pt(size)
    return p


def body_paragraph(doc, text, double=True):
    p = doc.add_paragraph()
    if double:
        p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(0)
    add_runs(p, fill(text))
    return p


# ---------------------------------------------------------------------------
# Table builders (return list-of-rows; header first)
# ---------------------------------------------------------------------------
def table1_data():
    return [
        ["Method", "Target level", "Core metric", "Primary input",
         "What it verifies", "Ref."],
        ["LINKO\n(Latent Information Normalization for Key Outcomes)",
         "Between studies",
         "Information contribution ratio (ICR)",
         "Published RCT summaries or IPD",
         "Which studies/endpoints carry the informative signal", "1"],
        ["IONE\n(Incoherence-Oriented Neutralisation and Extraction)",
         "Within studies",
         "Incoherence indicator C1 (= 1 \u2212 I\u00b2); within-stratum "
         "homogeneity W",
         "IPD or observational cohort",
         "Whether a population hides coherent subgroups (effect modification)",
         "2"],
        ["KOTHA\n(Knowledge-driven Observational-Trial Harmonisation Approach)",
         "Between study types",
         "Counterfactual power; harmonized estimate; OIS/TSA",
         "RCT + observational effect estimates; IPD",
         "How to harmonize randomized and observational evidence; "
         "\u201cno effect\u201d vs \u201cno information\u201d", "3"],
    ]


def table2_data():
    return [
        ["Interface", "Upstream output", "Downstream input", "Transformation"],
        ["IONE \u2192 LINKO", "Coherent subgroup labels", "Subgroup partition",
         "Recompute the ICR within each subgroup"],
        ["IONE \u2192 KOTHA (Module K)", "Subgroup risk profiles; C1",
         "Counterfactual power inputs",
         "Risk profiles parameterize per-stratum power simulation"],
        ["LINKO \u2192 KOTHA (Module T)", "Per-study / per-subgroup ICR",
         "Bayesian prior weight",
         "ICR scales the information borrowed in hierarchical integration"],
        ["LINKO, IONE \u2192 KOTHA (Module H)", "ICR; C1", "Reporting items",
         "Reported alongside OIS and TSA in the GRADE-linked read-out"],
    ]


def table3_data():
    rows = [["Stratum", "n", "14-day mortality", "Aspirin OR (95% CI)",
             "ICR (PCA)", "Power at pooled OR", "Additional patients to OIS"]]
    add_by = {d["stratum"]: d["additional_n"]
              for d in p3["additional_n_by_stratum"]}
    for i in range(4):
        s = strata[i]
        lo, hi = ci_from_log(s["logOR"], s["se"])
        rows.append([
            f"{i}", f"{s['n']:,}", pct(s["event_rate"]),
            f"{math.exp(s['logOR']):.2f} ({lo:.2f}\u2013{hi:.2f})",
            sci(p1[i]["icr_pca_reg"]),
            pct(p3["table"][i]["power_at_pooled_OR"], 0),
            f"{add_by[i]:,.0f}",
        ])
    rows.append([
        "Overall", f"{R['dataset']['N']:,}", pct(R["dataset"]["event_rate"]),
        f"{ov['OR']:.2f} ({ov_lo:.2f}\u2013{ov_hi:.2f})", "\u2014",
        f"{VALUES['INFOFRAC']} info. fraction",
        f"{p3['additional_n_overall']:,.0f}",
    ])
    return rows


def shade_cell(cell, hex_fill):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:fill"), hex_fill)
    tcPr.append(shd)


def render_table(doc, rows, title, footnote, font_size=9):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(title); r.bold = True
    r.font.name = FONT; r.font.size = Pt(10)
    tbl = doc.add_table(rows=len(rows), cols=len(rows[0]))
    tbl.style = "Table Grid"
    tbl.autofit = True
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            para = cell.paragraphs[0]
            para.paragraph_format.space_after = Pt(0)
            run = para.add_run(val)
            run.font.name = FONT
            run.font.size = Pt(font_size)
            if i == 0:
                run.bold = True
                shade_cell(cell, "D9E1F2")
    fn = doc.add_paragraph()
    fn.paragraph_format.space_before = Pt(3)
    r = fn.add_run(footnote)
    r.italic = True
    r.font.name = FONT
    r.font.size = Pt(8)
    return tbl


# ---------------------------------------------------------------------------
# Figure file map
# ---------------------------------------------------------------------------
FIG_SOURCES = {
    1: os.path.join(INTEG_FIGDIR, "schematic_overview.png"),
    2: os.path.join(INTEG_FIGDIR, "schematic_pipeline.png"),
    3: os.path.join(INTEG_FIGDIR, "pattern1_linko_ione.png"),
    4: os.path.join(INTEG_FIGDIR, "pattern2_linko_kotha.png"),
    5: os.path.join(INTEG_FIGDIR, "pattern3_ione_kotha.png"),
    6: os.path.join(INTEG_FIGDIR, "pattern4_onishi_full.png"),
}


def export_figures():
    os.makedirs(FIGOUT, exist_ok=True)
    mapping = {}
    for n, src in FIG_SOURCES.items():
        base = f"Figure_{n}"
        png = os.path.join(FIGOUT, base + ".png")
        tif = os.path.join(FIGOUT, base + ".tif")
        shutil.copyfile(src, png)
        im = Image.open(src).convert("RGB")
        im.save(tif, compression="tiff_lzw", dpi=(300, 300))
        mapping[n] = png
    return mapping


def add_image(doc, path, width_in=6.3):
    im = Image.open(path)
    w, h = im.size
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(10)
    p.add_run().add_picture(path, width=Inches(width_in))
    return p


def caption(doc, text, size=9):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(8)
    r = p.add_run(fill(text))
    r.font.name = FONT
    r.font.size = Pt(size)
    r.italic = True
    return p


# ---------------------------------------------------------------------------
# Word count (body only)
# ---------------------------------------------------------------------------
def body_word_count():
    total = 0
    for kind, val in C.BODY:
        if kind == "P":
            txt = CIT_RE.sub("", fill(val))
            total += len(txt.split())
    return total


# ---------------------------------------------------------------------------
# Main manuscript (submission format: no embedded figures)
# ---------------------------------------------------------------------------
def build_main(figmap):
    doc = Document()
    set_base_style(doc)
    sec = doc.sections[0]
    sec.top_margin = sec.bottom_margin = Inches(1)
    sec.left_margin = sec.right_margin = Inches(1)
    add_line_numbers(sec)
    add_page_numbers(sec)

    # Title page
    def center(text, size=12, bold=False, italic=False, after=6):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(after)
        r = p.add_run(text); r.bold = bold; r.italic = italic
        r.font.name = FONT; r.font.size = Pt(size)
        return p

    # Anonymized manuscript header (no author identity; the title page is a
    # separate file for AJE double-anonymous review).
    center("American Journal of Epidemiology", 11, italic=True, after=2)
    center(f"Article type: {C.ARTICLE_TYPE}", 11, after=12)
    center(C.TITLE, 15, bold=True, after=12)
    center("[Author names and affiliations removed for double-anonymous review; "
           "see separate title page.]", 9, italic=True, after=10)

    def kv(label, value, italic_v=False):
        p = doc.add_paragraph(); p.paragraph_format.space_after = Pt(2)
        r = p.add_run(label + " "); r.bold = True
        r.font.name = FONT; r.font.size = Pt(10)
        add_runs(p, value, size=10, italic=italic_v)
        return p

    doc.add_paragraph()
    kv("Running head:", C.SHORT_TITLE)
    kv("Abstract word count:", str(len(C.ABSTRACT.split())))
    kv("Main text word count:", str(body_word_count()))
    kv("Number of figures:", "6")
    kv("Number of tables:", "3")
    kv("Keywords:", "; ".join(C.KEYWORDS) + ".")
    kv("Data availability:", C.DATA_AVAILABILITY)
    kv("Use of artificial intelligence:", C.AI_USE)

    doc.add_page_break()

    # Abstract page
    heading(doc, "ABSTRACT")
    body_paragraph(doc, C.ABSTRACT)
    p = doc.add_paragraph(); p.paragraph_format.space_before = Pt(8)
    r = p.add_run("Keywords: "); r.bold = True; r.font.name = FONT; r.font.size = Pt(11)
    add_runs(p, "; ".join(C.KEYWORDS) + ".")

    # Highlights
    heading(doc, "HIGHLIGHTS")
    for h in C.HIGHLIGHTS:
        b = doc.add_paragraph(style="List Bullet")
        add_runs(b, fill(h))

    doc.add_page_break()

    # Body
    for kind, val in C.BODY:
        if kind == "H1":
            heading(doc, val)
        elif kind == "P":
            body_paragraph(doc, val)
        # FIG / TBL markers are ignored in submission format

    # References
    heading(doc, "REFERENCES")
    for i, ref in enumerate(C.REFERENCES, 1):
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = Inches(0.3)
        p.paragraph_format.first_line_indent = Inches(-0.3)
        p.paragraph_format.space_after = Pt(4)
        r = p.add_run(f"{i}. "); r.font.name = FONT; r.font.size = Pt(10)
        add_runs(p, ref, size=10)

    # Tables (each on its own page)
    doc.add_page_break()
    heading(doc, "TABLES")
    render_table(doc, table1_data(), C.TABLE_TITLES[1], C.TABLE_FOOTNOTES[1])
    doc.add_page_break()
    render_table(doc, table2_data(), C.TABLE_TITLES[2], C.TABLE_FOOTNOTES[2])
    doc.add_page_break()
    render_table(doc, table3_data(), C.TABLE_TITLES[3], C.TABLE_FOOTNOTES[3])

    # Figure legends
    doc.add_page_break()
    heading(doc, "FIGURE LEGENDS")
    for n in (1, 2, 3, 4, 5, 6):
        p = doc.add_paragraph(); p.paragraph_format.space_after = Pt(2)
        add_runs(p, fill(C.FIGURE_LEGENDS[n]))
        pa = doc.add_paragraph(); pa.paragraph_format.space_after = Pt(8)
        add_runs(pa, "Alt text: " + C.FIGURE_ALT_TEXT[n], italic=True)

    out = os.path.join(BUILD, "ONISHI_manuscript.docx")
    doc.save(out)
    return out


# ---------------------------------------------------------------------------
# Inline review copy (figures + tables placed at first mention)
# ---------------------------------------------------------------------------
def build_inline(figmap):
    doc = Document()
    set_base_style(doc)
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(0.9)
        s.left_margin = s.right_margin = Inches(0.9)

    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run(C.TITLE); r.bold = True; r.font.name = FONT; r.font.size = Pt(15)
    sub = doc.add_paragraph(); sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rr = sub.add_run("Review copy with figures and tables placed inline "
                     "(not the submission-format file).")
    rr.italic = True; rr.font.name = FONT; rr.font.size = Pt(10)

    heading(doc, "ABSTRACT")
    body_paragraph(doc, C.ABSTRACT, double=False)

    for kind, val in C.BODY:
        if kind == "H1":
            heading(doc, val)
        elif kind == "P":
            body_paragraph(doc, val, double=False)
        elif kind == "FIG":
            add_image(doc, figmap[val])
            caption(doc, C.FIGURE_LEGENDS[val])
        elif kind == "TBL":
            data = {1: table1_data, 2: table2_data, 3: table3_data}[val]()
            render_table(doc, data, C.TABLE_TITLES[val], C.TABLE_FOOTNOTES[val])

    heading(doc, "REFERENCES")
    for i, ref in enumerate(C.REFERENCES, 1):
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = Inches(0.3)
        p.paragraph_format.first_line_indent = Inches(-0.3)
        p.paragraph_format.space_after = Pt(3)
        r = p.add_run(f"{i}. "); r.font.name = FONT; r.font.size = Pt(10)
        add_runs(p, ref, size=10)

    out = os.path.join(BUILD, "ONISHI_manuscript_inline.docx")
    doc.save(out)
    return out


def build_tables_doc():
    doc = Document()
    set_base_style(doc)
    t = doc.add_paragraph(); r = t.add_run("ONISHI \u2014 Tables (editable)")
    r.bold = True; r.font.name = FONT; r.font.size = Pt(14)
    render_table(doc, table1_data(), C.TABLE_TITLES[1], C.TABLE_FOOTNOTES[1])
    doc.add_paragraph()
    render_table(doc, table2_data(), C.TABLE_TITLES[2], C.TABLE_FOOTNOTES[2])
    doc.add_paragraph()
    render_table(doc, table3_data(), C.TABLE_TITLES[3], C.TABLE_FOOTNOTES[3])
    out = os.path.join(BUILD, "ONISHI_tables.docx")
    doc.save(out)
    return out


# ---------------------------------------------------------------------------
# Editable figures pptx (one per slide)
# ---------------------------------------------------------------------------
def build_pptx(figmap):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    items = [(n, f"Figure {n}", C.FIGURE_LEGENDS[n]) for n in (1, 2, 3, 4, 5, 6)]
    for key, title, legend in items:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(PInches(0.4), PInches(0.15),
                                      PInches(12.5), PInches(0.6))
        tf = tb.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; para.text = title
        para.font.size = PPt(22); para.font.bold = True
        para.alignment = PP_ALIGN.LEFT

        path = figmap[key]
        im = Image.open(path); w, h = im.size
        avail_w, avail_h = 12.5, 5.4
        ratio = min(avail_w / (w / 96.0), avail_h / (h / 96.0))
        pic_w = (w / 96.0) * ratio
        pic_h = (h / 96.0) * ratio
        left = (13.333 - pic_w) / 2
        slide.shapes.add_picture(path, PInches(left), PInches(0.9),
                                 height=PInches(pic_h))

        cb = slide.shapes.add_textbox(PInches(0.4), PInches(6.5),
                                      PInches(12.5), PInches(0.9))
        cf = cb.text_frame; cf.word_wrap = True
        cpara = cf.paragraphs[0]; cpara.text = fill(legend)
        cpara.font.size = PPt(10)
    out = os.path.join(BUILD, "ONISHI_figures.pptx")
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# Cover letter
# ---------------------------------------------------------------------------
def build_cover_letter():
    doc = Document()
    set_base_style(doc)
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(1)
        s.left_margin = s.right_margin = Inches(1)

    def para(text, bold=False, after=8, size=11):
        p = doc.add_paragraph(); p.paragraph_format.space_after = Pt(after)
        p.paragraph_format.line_spacing = 1.15
        r = p.add_run(text); r.bold = bold
        r.font.name = FONT; r.font.size = Pt(size)
        return p

    today = datetime.date.today().strftime("%B %-d, %Y")
    para(today)
    para("The Editor-in-Chief")
    para("American Journal of Epidemiology")
    doc.add_paragraph()
    para("Dear Editor,")
    para("Please consider our manuscript, \u201c" + C.TITLE + ",\u201d for "
         "publication in the American Journal of Epidemiology as a Practice of "
         "Epidemiology article.")
    para("Evidence synthesis can be compromised at three distinct levels\u2014"
         "between studies (unequal information content), within studies (hidden "
         "coherent subgroups), and between study types (harmonizing randomized "
         "and observational evidence, and distinguishing \u201cno effect\u201d "
         "from \u201cno information\u201d). We present ONISHI, a framework that "
         "pairs each level with one dedicated method (LINKO, IONE, and KOTHA, "
         "respectively) and chains them into a single pipeline on shared data, "
         "with explicit hand-off interfaces between methods. We illustrate the "
         "whole pipeline on public individual patient data from the International "
         "Stroke Trial, showing how an inconclusive result can be characterized "
         "honestly as coherent but information-limited, with the additional "
         "evidence required made explicit.")
    para("Relationship to our other work. The three component methods are "
         "described in separate manuscripts currently under review (cited as "
         "preprints in the reference list). The present paper does not restate "
         "those methods; its contribution is the cross-level integration: the "
         "conceptual mapping of three methods to three levels, the analytic "
         "capabilities that arise from combining them, the defined data hand-offs "
         "between methods, and a worked end-to-end application. We note this to be "
         "transparent about overlap and to assist reviewer selection.")
    para("Use of artificial intelligence. In the interest of full transparency, "
         "we disclose that an AI assistant was used to help draft and format the "
         "manuscript and to implement the analysis code. All numbers were "
         "regenerated from the public dataset and verified against the source "
         "data; the authors reviewed and approved every claim and take full "
         "responsibility for the content.")
    para("Declarations. This manuscript is original, is not under consideration "
         "elsewhere, and has not been published previously. All authors have "
         "approved the submission and declare no conflicts of interest. The "
         "illustrative analysis uses a fully anonymized, publicly available "
         "dataset and did not require ethical approval. Suggested and "
         "non-preferred reviewers are provided in a separate document.")
    para("Thank you for considering our work.")
    doc.add_paragraph()
    para("Sincerely,")
    para("[Author One], on behalf of all authors")
    para("[Affiliation, email]")
    out = os.path.join(BUILD, "ONISHI_cover_letter.docx")
    doc.save(out)
    return out


def build_highlights():
    doc = Document(); set_base_style(doc)
    r = doc.add_paragraph().add_run("Highlights"); r.bold = True
    r.font.name = FONT; r.font.size = Pt(14)
    for h in C.HIGHLIGHTS:
        b = doc.add_paragraph(style="List Bullet"); add_runs(b, fill(h))
    doc.add_paragraph()
    r = doc.add_paragraph().add_run("Keywords"); r.bold = True
    r.font.name = FONT; r.font.size = Pt(14)
    p = doc.add_paragraph(); add_runs(p, "; ".join(C.KEYWORDS) + ".")
    out = os.path.join(BUILD, "ONISHI_highlights_keywords.docx")
    doc.save(out)
    return out


def build_reviewers():
    doc = Document(); set_base_style(doc)
    r = doc.add_paragraph().add_run("Suggested and non-preferred reviewers")
    r.bold = True; r.font.name = FONT; r.font.size = Pt(14)

    p = doc.add_paragraph()
    add_runs(p, "The authors suggest the following reviewers on the basis of "
                "relevant methodological expertise in evidence synthesis, "
                "treatment-effect heterogeneity, and evidence appraisal. Contact "
                "details should be completed and current conflicts confirmed by "
                "the authors before submission; none is a co-author of this work.")

    r = doc.add_paragraph().add_run("Preferred reviewers"); r.bold = True
    r.font.name = FONT; r.font.size = Pt(12)
    preferred = [
        ("Georgia Salanti, PhD", "Institute of Social and Preventive Medicine, "
         "University of Bern, Switzerland",
         "Evidence synthesis and network meta-analysis methodology."),
        ("Orestis Efthimiou, PhD", "University of Bern, Switzerland",
         "Meta-analysis methods and Bayesian evidence synthesis."),
        ("Issa J. Dahabreh, MD, ScD", "Harvard T.H. Chan School of Public "
         "Health, USA",
         "Transportability and integration of randomized and observational "
         "evidence; target-trial thinking."),
        ("Ian R. White, PhD", "MRC Clinical Trials Unit, University College "
         "London, UK",
         "Statistical methods for meta-analysis and individual patient data."),
        ("David M. Kent, MD, MS", "Tufts Medical Center, USA",
         "Predictive approaches to treatment-effect heterogeneity (PATH)."),
    ]
    for name, aff, why in preferred:
        b = doc.add_paragraph(style="List Bullet")
        run = b.add_run(name + " \u2014 "); run.bold = True
        run.font.name = FONT; run.font.size = Pt(11)
        add_runs(b, aff + ". " + why + " [email to be added]")

    r = doc.add_paragraph().add_run("Non-preferred (opposed) reviewers"); r.bold = True
    r.font.name = FONT; r.font.size = Pt(12)
    p = doc.add_paragraph()
    add_runs(p, "Because authorship is not yet finalized in this version, "
                "specific opposed reviewers cannot be confirmed here. The authors "
                "should exclude individuals with a genuine conflict, applying the "
                "following criteria, and list names before submission:")
    for g in ["Close collaborators or co-authors within the past three years.",
              "Individuals at the authors' own institutions.",
              "Direct competitors developing overlapping proprietary frameworks "
              "who may have a conflict of interest.",
              "[Named individual 1 \u2014 reason]",
              "[Named individual 2 \u2014 reason]"]:
        b = doc.add_paragraph(style="List Bullet"); add_runs(b, g)
    out = os.path.join(BUILD, "ONISHI_reviewer_suggestions.docx")
    doc.save(out)
    return out


def build_strobe():
    doc = Document(); set_base_style(doc)
    r = doc.add_paragraph().add_run(
        "STROBE Statement \u2014 checklist for the illustrative analysis")
    r.bold = True; r.font.name = FONT; r.font.size = Pt(13)
    p = doc.add_paragraph()
    add_runs(p, "The illustrative application (\u201cIllustrative "
                "Application\u201d) is a secondary analysis of publicly available "
                "individual patient data and is reported following the "
                "Strengthening the Reporting of Observational Studies in "
                "Epidemiology (STROBE) recommendations. Item locations refer to "
                "the manuscript sections. STROBE is a reporting guide; the "
                "methodological sections of the paper are not observational "
                "study reports.", size=10)
    items = [
        ("1", "Title and abstract",
         "Design indicated; abstract summarizes design, data, and findings.",
         "Title; Abstract"),
        ("2", "Background/rationale", "Rationale for the framework and the "
         "worked example.", "Introduction"),
        ("3", "Objectives", "Objective: demonstrate the integrated pipeline on "
         "shared data.", "Introduction; Illustrative Application"),
        ("4", "Study design", "Secondary analysis of a randomized trial's public "
         "IPD, with latent risk stratification.", "Illustrative Application"),
        ("5", "Setting", "International Stroke Trial, 1991\u20131996; public "
         "database.", "Illustrative Application; Refs 12, 13"),
        ("6", "Participants", "18,451 patients with acute ischaemic stroke; "
         "eligibility as in the source trial.", "Illustrative Application"),
        ("7", "Variables", "Treatment: randomized aspirin allocation; outcome: "
         "14-day death; 25 baseline covariates.", "Illustrative Application"),
        ("8", "Data sources/measurement", "Variables as recorded in the public "
         "IST database.", "Illustrative Application; Ref 12"),
        ("9", "Bias", "Randomized treatment allocation limits confounding; "
         "stratification is exploratory (acknowledged).", "Discussion"),
        ("10", "Study size", "Whole available sample analysed; information-size "
         "projections reported.", "Illustrative Application; Table 3"),
        ("11", "Quantitative variables", "Latent risk strata defined by predicted "
         "outcome probability (four groups).", "Illustrative Application"),
        ("12", "Statistical methods", "IONE stratification; PCA-based ICR; "
         "random-effects and ICR-weighted pooling; counterfactual power; OIS/TSA; "
         "Bayesian power-prior integration.", "Component Methods; Illustrative "
         "Application"),
        ("13", "Participants (results)", "Sample and strata sizes reported.",
         "Table 3"),
        ("14", "Descriptive data", "Event rates by stratum reported.",
         "Illustrative Application; Table 3"),
        ("15", "Outcome data", "14-day deaths and odds ratios reported.",
         "Illustrative Application; Table 3; Figures 3\u20136"),
        ("16", "Main results", "Odds ratios with 95% CIs; ICR; power; "
         "information fraction; additional evidence required.",
         "Illustrative Application; Table 3; Figures 3\u20136"),
        ("17", "Other analyses", "Between-country pooling; sensitivity to "
         "weighting scheme.", "Illustrative Application; Figure 4"),
        ("18", "Key results", "Coherent effect but information-limited evidence.",
         "Illustrative Application; Discussion"),
        ("19", "Limitations", "Exploratory stratification; IPD dependence; prior "
         "sensitivity; single illustrative case.", "Discussion"),
        ("20", "Interpretation", "Reframes an inconclusive trial as interim; "
         "quantifies evidence needed.", "Discussion"),
        ("21", "Generalisability", "Framework general; example specific to the "
         "IST setting.", "Discussion"),
        ("22", "Funding", "Stated on the title page.", "Title page"),
    ]
    tbl = doc.add_table(rows=len(items) + 1, cols=4)
    tbl.style = "Table Grid"
    hdr = ["Item", "STROBE recommendation", "How addressed", "Location"]
    for j, h in enumerate(hdr):
        c = tbl.cell(0, j); c.text = ""
        rr = c.paragraphs[0].add_run(h); rr.bold = True
        rr.font.name = FONT; rr.font.size = Pt(9)
        shade_cell(c, "D9E1F2")
    for i, (num, rec, how, loc) in enumerate(items, 1):
        for j, val in enumerate((num, rec, how, loc)):
            c = tbl.cell(i, j); c.text = ""
            rr = c.paragraphs[0].add_run(val)
            rr.font.name = FONT; rr.font.size = Pt(8.5)
    out = os.path.join(BUILD, "ONISHI_STROBE_checklist.docx")
    doc.save(out)
    return out


def build_title_page():
    """Separate title page (the only identity-revealing file), AJE form order."""
    doc = Document()
    set_base_style(doc)
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(1)
        s.left_margin = s.right_margin = Inches(1)

    h = doc.add_paragraph()
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = h.add_run("American Journal of Epidemiology Submitted Manuscript")
    r.bold = True; r.font.name = FONT; r.font.size = Pt(12)
    h.paragraph_format.space_after = Pt(12)

    authors = ", ".join(a["name"] + ("*" if a["corr"] else "")
                        for a in C.AUTHORS)
    affs = "; ".join(f"{n}. {a}" for n, a in C.AFFILIATIONS.items())
    rows = [
        ("Special Collection:", C.SPECIAL_COLLECTION),
        ("Title:", C.TITLE),
        ("Authors:", authors + "  (* corresponding author)"),
        ("ORCiD IDs:", C.ORCID_IDS),
        ("Correspondence Address:", C.CORR_ADDRESS),
        ("Joint Authorship:", C.JOINT_AUTHORSHIP),
        ("Affiliations:", affs),
        ("Key words:", "; ".join(C.KEYWORDS) + "."),
        ("Acknowledgments:", C.ACKNOWLEDGEMENTS),
        ("Funding:", C.FUNDING),
        ("Conflict of Interest:", C.CONFLICTS),
        ("Disclaimer:", C.DISCLAIMER),
        ("Data Availability Statement:", C.DATA_AVAILABILITY),
    ]
    for label, value in rows:
        p = doc.add_paragraph(); p.paragraph_format.space_after = Pt(6)
        rr = p.add_run(label + " "); rr.bold = True
        rr.font.name = FONT; rr.font.size = Pt(11)
        add_runs(p, value, size=11)

    out = os.path.join(BUILD, "ONISHI_title_page.docx")
    doc.save(out)
    return out


def zip_all(files):
    zpath = os.path.join(BUILD, "ONISHI_submission_package.zip")
    with zipfile.ZipFile(zpath, "w", zipfile.ZIP_DEFLATED) as z:
        for f in files:
            z.write(f, os.path.join("ONISHI_submission", os.path.basename(f)))
        for fn in sorted(os.listdir(FIGOUT)):
            z.write(os.path.join(FIGOUT, fn),
                    os.path.join("ONISHI_submission", "figures", fn))
    return zpath


def main():
    if os.path.isdir(BUILD):
        shutil.rmtree(BUILD)
    os.makedirs(BUILD)
    figmap = export_figures()
    files = []
    files.append(build_title_page())
    files.append(build_main(figmap))
    files.append(build_inline(figmap))
    files.append(build_tables_doc())
    files.append(build_pptx(figmap))
    files.append(build_cover_letter())
    files.append(build_highlights())
    files.append(build_reviewers())
    files.append(build_strobe())
    zpath = zip_all(files)
    print("Body word count:", body_word_count())
    print("Abstract word count:", len(C.ABSTRACT.split()))
    print("Built:")
    for f in files:
        print("  ", os.path.relpath(f, HERE))
    print("   figures/ ->", len(os.listdir(FIGOUT)), "files")
    print("ZIP:", os.path.relpath(zpath, HERE))


if __name__ == "__main__":
    main()
